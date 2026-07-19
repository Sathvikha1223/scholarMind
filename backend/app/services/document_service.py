import os
import uuid
import logging
from pathlib import Path
from sqlalchemy.orm import Session
from fastapi import HTTPException, status, UploadFile
from app.models.document import Document, DocumentStatus
from app.core.config import settings

logger = logging.getLogger(__name__)

ALLOWED_EXTENSIONS = {
    "pdf", "docx", "pptx", "txt", "md",
    "py", "java", "js", "c", "cpp", "sql", "html", "css", "ts", "tsx"
}


def validate_file(file: UploadFile) -> str:
    ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"File type '.{ext}' not supported. Allowed: {', '.join(ALLOWED_EXTENSIONS)}",
        )
    return ext


async def save_uploaded_file(file: UploadFile, user_id: str, db: Session) -> Document:
    ext = validate_file(file)
    content = await file.read()
    if len(content) > settings.max_file_size_bytes:
        raise HTTPException(
            status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
            detail=f"File exceeds {settings.MAX_FILE_SIZE_MB}MB limit",
        )

    upload_dir = Path(settings.UPLOAD_DIR) / user_id
    upload_dir.mkdir(parents=True, exist_ok=True)

    stored_filename = f"{uuid.uuid4()}.{ext}"
    file_path = upload_dir / stored_filename

    with open(file_path, "wb") as f:
        f.write(content)

    doc = Document(
        user_id=user_id,
        filename=stored_filename,
        original_filename=file.filename,
        file_type=ext,
        file_size=len(content),
        status=DocumentStatus.pending.value,
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)
    return doc


def extract_text_from_document(doc: Document) -> tuple[str, int]:
    """Extract raw text and page count from document file. Returns (text, page_count)."""
    file_path = Path(settings.UPLOAD_DIR) / doc.user_id / doc.filename
    ext = doc.file_type.lower()

    try:
        if ext == "pdf":
            return _extract_pdf(file_path)
        elif ext == "docx":
            return _extract_docx(file_path)
        elif ext == "pptx":
            return _extract_pptx(file_path)
        else:
            # Plain text formats
            text = file_path.read_text(encoding="utf-8", errors="ignore")
            return text, 1
    except Exception as e:
        logger.error(f"Text extraction failed for {doc.filename}: {e}")
        raise


def _extract_pdf(file_path: Path) -> tuple[str, int]:
    import fitz  # PyMuPDF
    doc = fitz.open(str(file_path))
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text")
        if text.strip():
            pages.append(f"[Page {i + 1}]\n{text}")
    doc.close()
    return "\n\n".join(pages), len(doc)


def _extract_docx(file_path: Path) -> tuple[str, int]:
    from docx import Document as DocxDocument
    doc = DocxDocument(str(file_path))
    text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    return text, 1


def _extract_pptx(file_path: Path) -> tuple[str, int]:
    from pptx import Presentation
    prs = Presentation(str(file_path))
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)
        if slide_content:
            slides_text.append(f"[Slide {i + 1}]\n" + "\n".join(slide_content))
    return "\n\n".join(slides_text), len(prs.slides)


def get_user_documents(db: Session, user_id: str) -> list[Document]:
    return db.query(Document).filter(Document.user_id == user_id).order_by(Document.created_at.desc()).all()


def get_document_by_id(db: Session, doc_id: str, user_id: str) -> Document:
    doc = db.query(Document).filter(Document.id == doc_id, Document.user_id == user_id).first()
    if not doc:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Document not found")
    return doc


def delete_document(db: Session, doc_id: str, user_id: str) -> None:
    doc = get_document_by_id(db, doc_id, user_id)
    file_path = Path(settings.UPLOAD_DIR) / user_id / doc.filename
    if file_path.exists():
        file_path.unlink()
    db.delete(doc)
    db.commit()
